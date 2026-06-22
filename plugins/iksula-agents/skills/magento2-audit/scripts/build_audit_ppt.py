#!/usr/bin/env python3
import sys, json, os, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

INK=RGBColor(0x15,0x15,0x18); RED=RGBColor(0x9A,0x0D,0x15); WHITE=RGBColor(0xFF,0xFF,0xFF)
LGRAY=RGBColor(0xF4,0xF4,0xF6); MGRAY=RGBColor(0x5E,0x5E,0x66); HAIR=RGBColor(0xDD,0xDD,0xE2)
AMBER=RGBColor(0xD9,0x8A,0x1F); GREEN=RGBColor(0x2E,0x7D,0x32); FONT="Carlito"

def score_color(s):
    try: s=float(s)
    except: return MGRAY
    return RED if s<=2 else (AMBER if s<3.5 else GREEN)
def rating_color(r): return {"good":GREEN,"avg":AMBER,"poor":RED}.get(str(r).strip().lower(),MGRAY)
def pcolor(v): return {"high":RED,"medium":AMBER,"low":MGRAY}.get(str(v).strip().lower(),MGRAY)
def lat_color(v):
    m=re.search(r'([\d.]+)',str(v)); s='s' in str(v).lower()
    if not m: return MGRAY
    x=float(m.group(1))
    if s and x>=10: return RED
    if s and x>=3: return AMBER
    return GREEN
def apdex_color(v):
    try: x=float(re.search(r'([\d.]+)',str(v)).group(1))
    except: return MGRAY
    return RED if x<0.5 else (AMBER if x<0.85 else GREEN)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height; BLANK=prs.slide_layouts[6]
data=json.load(open(sys.argv[1],encoding="utf-8")); meta=data.get("meta",{}); dims=data.get("dimensions",[])
out=sys.argv[2] if len(sys.argv)>2 else os.path.splitext(sys.argv[1])[0]+".pptx"
CLIENT=meta.get("client","Audit"); PAGE=[0]
ML=Inches(0.62); CW=Inches(12.1)   # left margin, content width

def slide(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,fill=None,line=None,shape=MSO_SHAPE.RECTANGLE):
    sp=s.shapes.add_shape(shape,x,y,w,h); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    return sp
def tbox(s,x,y,w,h,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0; return tf
def _nb(p): p._p.get_or_add_pPr().append(p._p.get_or_add_pPr().makeelement(qn('a:buNone'),{}))
def _bu(p,c):
    pPr=p._p.get_or_add_pPr()
    pPr.insert(0,pPr.makeelement(qn('a:buChar'),{'char':'•'}))
    pPr.insert(0,pPr.makeelement(qn('a:buFont'),{'typeface':'Carlito'}))
    cl=pPr.makeelement(qn('a:buClr'),{}); cl.append(cl.makeelement(qn('a:srgbClr'),{'val':'%02X%02X%02X'%(c[0],c[1],c[2])})); pPr.insert(0,cl)
def para(tf,text,size,color=INK,bold=False,first=False,align=PP_ALIGN.LEFT,sa=6,italic=False,bullet=False,spc=None):
    p=tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment=align; p.space_after=Pt(sa); p.space_before=Pt(0)
    if spc is not None: p.line_spacing=spc
    r=p.add_run(); r.text=str(text); f=r.font; f.name=FONT; f.size=Pt(size); f.bold=bold; f.italic=italic; f.color.rgb=color
    _bu(p,color) if bullet else _nb(p); return p
def chip(s,x,y,w,text,fill,fg=WHITE,size=10):
    c=rect(s,x,y,w,Inches(0.32),fill=fill,shape=MSO_SHAPE.ROUNDED_RECTANGLE); c.adjustments[0]=0.5
    tf=c.text_frame; tf.word_wrap=False
    tf.margin_left=Pt(7);tf.margin_right=Pt(7);tf.margin_top=Pt(0);tf.margin_bottom=Pt(0)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=str(text); r.font.name=FONT; r.font.size=Pt(size); r.font.bold=True; r.font.color.rgb=fg
    return c
def footer(s):
    PAGE[0]+=1
    rect(s,ML,Inches(7.06),CW,Pt(0.75),fill=HAIR)
    para(tbox(s,ML,Inches(7.12),Inches(9),Inches(0.3)),CLIENT+"   |   E-commerce Audit   |   Confidential",8.5,MGRAY,first=True)
    para(tbox(s,Inches(10.0),Inches(7.12),Inches(2.73),Inches(0.3)),str(PAGE[0]),8.5,MGRAY,first=True,align=PP_ALIGN.RIGHT)
def H(s,kicker,title,sub=None):
    rect(s,0,0,SW,Inches(1.28),fill=INK); rect(s,0,Inches(1.28),SW,Inches(0.05),fill=RED)
    para(tbox(s,ML,Inches(0.26),CW,Inches(0.3)),kicker.upper(),12,RED,bold=True,first=True)
    para(tbox(s,ML,Inches(0.57),CW,Inches(0.62)),title,26,WHITE,bold=True,first=True)
    footer(s)
    if sub: para(tbox(s,ML,Inches(1.46),CW,Inches(0.6)),sub,12,MGRAY,first=True)
    return Inches(1.46) if not sub else Inches(2.05)

# ---------- table helper (striped, padded, readable) ----------
def table(s,x,y,w,colw,rows,header_h=0.46,row_h=0.4,fs=10.5,hfs=10):
    nC=len(colw); nR=len(rows)
    t=s.shapes.add_table(nR,nC,x,y,w,Inches(header_h+row_h*(nR-1))).table
    t._tbl.tblPr.set('firstRow','0'); t._tbl.tblPr.set('bandRow','0')
    for i,cw in enumerate(colw): t.columns[i].width=cw
    t.rows[0].height=Inches(header_h)
    for ri in range(1,nR): t.rows[ri].height=Inches(row_h)
    return t
def cell(c,blocks,fill,va=MSO_ANCHOR.MIDDLE):
    c.margin_left=Inches(0.1);c.margin_right=Inches(0.08);c.margin_top=Inches(0.05);c.margin_bottom=Inches(0.05)
    c.vertical_anchor=va; c.fill.solid(); c.fill.fore_color.rgb=fill
    tf=c.text_frame; tf.word_wrap=True
    for i,b in enumerate(blocks):
        txt,sz,col,bold,al=b[0],b[1],b[2],b[3],(b[4] if len(b)>4 else PP_ALIGN.LEFT)
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=al; p.space_after=Pt(2); p.space_before=Pt(0)
        r=p.add_run(); r.text=str(txt); f=r.font; f.name=FONT; f.size=Pt(sz); f.bold=bold; f.color.rgb=col
        _nb(p)

# ============================== 1. COVER ==============================
s=slide(); rect(s,0,0,SW,SH,fill=INK)
rect(s,ML,Inches(2.05),Inches(0.85),Inches(0.09),fill=RED)
para(tbox(s,ML,Inches(0.85),CW,Inches(0.4)),"E-COMMERCE PLATFORM AUDIT",14,RED,bold=True,first=True)
para(tbox(s,ML,Inches(2.4),CW,Inches(1.3)),CLIENT,46,WHITE,bold=True,first=True,sa=2)
para(tbox(s,ML,Inches(3.55),CW,Inches(0.5)),meta.get("site_url",""),20,RGBColor(0xC8,0xC8,0xD0))
tf=tbox(s,ML,Inches(5.5),CW,Inches(1.4))
para(tf,"Platform:  "+meta.get("platform","-"),14,RGBColor(0xC8,0xC8,0xD0),first=True,sa=6)
para(tf,"Date:  "+meta.get("audit_date","")+("      Prepared by:  "+meta.get("auditor","") if meta.get("auditor") else ""),13,MGRAY)

# ============================== 2. COVERAGE ==============================
s=slide(); H(s,"01  Basis of audit","Coverage & confidence",
  "The audit depth matches the access provided. Items without access are marked Not assessed - never assumed.")
cov=data.get("coverage",[])
rows=[["ACCESS SOURCE","AVAILABLE","CONFIDENCE","NOTES"]]+[[c.get("source",""),"Yes" if c.get("available") else "No",c.get("confidence","-"),c.get("note","")] for c in cov]
t=table(s,ML,Inches(2.35),CW,[Inches(3.0),Inches(1.5),Inches(1.9),Inches(5.7)],rows,row_h=0.62)
for ri,row in enumerate(rows):
    for ci,val in enumerate(row):
        if ri==0: cell(t.cell(ri,ci),[(val,10,WHITE,True)],INK)
        else:
            fl=WHITE if ri%2 else LGRAY
            if ci==1: cell(t.cell(ri,ci),[(val,11,GREEN if val=="Yes" else RED,True,PP_ALIGN.CENTER)],fl)
            elif ci==2: cell(t.cell(ri,ci),[(val,11,score_color({"high":4,"medium":3,"low":1}.get(str(val).lower(),3)),True,PP_ALIGN.CENTER)],fl)
            elif ci==0: cell(t.cell(ri,ci),[(val,11,INK,True)],fl)
            else: cell(t.cell(ri,ci),[(val,9.5,MGRAY,False)],fl)

# ============================== 3. SCORECARD ==============================
s=slide(); H(s,"02  Executive summary","Maturity scorecard")
scs=[d.get("score") for d in dims if isinstance(d.get("score"),(int,float))]
overall=meta.get("overall_score") or (round(sum(scs)/len(scs),1) if scs else None)
summ=data.get("summary") or meta.get("summary")
if summ: para(tbox(s,ML,Inches(1.5),Inches(9.2),Inches(1.45)),summ,11.5,MGRAY,first=True,spc=1.15)
if overall is not None:
    rect(s,Inches(10.1),Inches(1.5),Inches(2.62),Inches(1.5),fill=INK)
    para(tbox(s,Inches(10.1),Inches(1.66),Inches(2.62),Inches(0.3)),"OVERALL MATURITY",10,RED,bold=True,first=True,align=PP_ALIGN.CENTER)
    para(tbox(s,Inches(10.1),Inches(1.98),Inches(2.62),Inches(0.9),MSO_ANCHOR.MIDDLE),str(overall)+" / 5",32,score_color(overall),bold=True,first=True,align=PP_ALIGN.CENTER)
by=Inches(3.35); bh=Inches(0.5); gap=Inches(0.28)
for i,d in enumerate(dims):
    y=by+i*(bh+gap)
    para(tbox(s,ML,y-Inches(0.04),Inches(3.4),bh,MSO_ANCHOR.MIDDLE),str(i+1)+".  "+d.get('name',''),12,INK,bold=True,first=True)
    sc=d.get("score",0) or 0
    rect(s,Inches(4.2),y+Inches(0.06),Inches(6.6),Inches(0.34),fill=LGRAY)
    try: fr=max(0.0,min(1.0,float(sc)/5.0))
    except: fr=0
    if fr>0: rect(s,Inches(4.2),y+Inches(0.06),Inches(6.6*fr),Inches(0.34),fill=score_color(sc))
    para(tbox(s,Inches(10.95),y-Inches(0.04),Inches(0.7),bh,MSO_ANCHOR.MIDDLE),str(sc)+"/5",12,score_color(sc),bold=True,first=True)
    if d.get("confidence"): para(tbox(s,Inches(11.6),y-Inches(0.04),Inches(1.1),bh,MSO_ANCHOR.MIDDLE),str(d["confidence"]),8.5,MGRAY,first=True)
# legend
ly=Inches(6.6)
for j,(lab,c) in enumerate([("1-2  Critical/Weak",RED),("3  Acceptable",AMBER),("4-5  Good/Strong",GREEN)]):
    rect(s,ML+j*Inches(2.6),ly,Inches(0.22),Inches(0.22),fill=c)
    para(tbox(s,ML+j*Inches(2.6)+Inches(0.3),ly-Inches(0.02),Inches(2.3),Inches(0.3)),lab,9.5,MGRAY,first=True)

# ============================== 4. PLATFORM HEALTH ==============================
ts=data.get("tech_snapshot")
if ts:
    s=slide(); H(s,"03  Platform health","Performance & security snapshot", ts.get("note"))
    para(tbox(s,ML,Inches(2.25),Inches(6),Inches(0.3)),"PERFORMANCE",12,INK,bold=True,first=True)
    perf=ts.get("performance",[]); cw=Inches(2.9); ch=Inches(1.0); gx=Inches(0.16)
    for i,m in enumerate(perf[:8]):
        r=i//2; c=i%2; x=ML+c*(cw+gx); y=Inches(2.68)+r*(ch+Inches(0.16)); col=rating_color(m.get("rating"))
        rect(s,x,y,cw,ch,fill=LGRAY); rect(s,x,y,Inches(0.09),ch,fill=col)
        para(tbox(s,x+Inches(0.28),y+Inches(0.12),cw-Inches(0.4),Inches(0.35)),m.get("label",""),10,MGRAY,first=True)
        para(tbox(s,x+Inches(0.28),y+Inches(0.42),cw-Inches(0.4),Inches(0.5)),m.get("value",""),19,col,bold=True,first=True)
    sx=Inches(7.05)
    para(tbox(s,sx,Inches(2.25),Inches(5.4),Inches(0.3)),"SECURITY",12,INK,bold=True,first=True)
    sec=ts.get("security",[])
    rows=[["CHECK","STATUS"]]+[[x.get("check",""),x.get("status","")] for x in sec]
    t=table(s,sx,Inches(2.68),Inches(5.5),[Inches(2.7),Inches(2.8)],rows,row_h=0.52,header_h=0.42)
    for ri,row in enumerate(rows):
        for ci,val in enumerate(row):
            if ri==0: cell(t.cell(ri,ci),[(val,10,WHITE,True)],INK)
            else:
                fl=WHITE if ri%2 else LGRAY
                if ci==1: cell(t.cell(ri,ci),[(val,9.5,rating_color(sec[ri-1].get("rating")),True)],fl)
                else: cell(t.cell(ri,ci),[(val,10,INK,True)],fl)

# ============================== 5/6/7 TOP-5 ==============================
def top5(num,kicker,title,items,accent=RED):
    s=slide(); H(s,kicker,title)
    if not items: para(tbox(s,ML,Inches(1.7),CW,Inches(1)),"(none recorded)",13,MGRAY,first=True,italic=True); return
    y0=Inches(1.75); rh=Inches(0.98)
    for i,it in enumerate(items[:5]):
        y=y0+i*rh
        lead=it.get("title") if isinstance(it,dict) else str(it)
        body=(it.get("impact") or it.get("detail") or "") if isinstance(it,dict) else ""
        rect(s,ML,y,Inches(0.62),Inches(0.62),fill=INK)
        para(tbox(s,ML,y,Inches(0.62),Inches(0.62),MSO_ANCHOR.MIDDLE),str(i+1),22,WHITE,bold=True,first=True,align=PP_ALIGN.CENTER)
        para(tbox(s,Inches(1.45),y+Inches(0.02),Inches(11.2),Inches(0.45)),lead,15,INK,bold=True,first=True)
        if body: para(tbox(s,Inches(1.45),y+Inches(0.46),Inches(11.2),Inches(0.45)),body,11.5,MGRAY,first=True)
top5(5,"04  Key risks","Top 5 risks",data.get("top_risks",[]))
top5(6,"05  Conversion","Top 5 conversion blockers",data.get("top_blockers",[]))
top5(7,"06  Growth","Top 5 modernization opportunities",data.get("top_opportunities",[]))

# ============================== 9. DIMENSION FINDINGS ==============================
HDR=[("FINDING & CURRENT STATE",),("RECOMMENDATION",),("PRIORITY / IMPACT / EFFORT",)]
CW3=[Inches(5.1),Inches(4.4),Inches(2.6)]
for i,d in enumerate(dims):
    fs=d.get("findings",[])
    if d.get("not_assessed") or not fs:
        s=slide(); H(s,"Dimension "+str(i+1),d.get("name",""))
        if d.get("confidence"): chip(s,Inches(11.0),Inches(0.62),Inches(1.7),"Confidence: "+str(d["confidence"]),MGRAY)
        rect(s,ML,Inches(2.7),CW,Inches(2.0),fill=LGRAY); rect(s,ML,Inches(2.7),Inches(0.1),Inches(2.0),fill=RED)
        tf=tbox(s,Inches(1.1),Inches(3.0),Inches(10.6),Inches(1.5))
        para(tf,"INSUFFICIENT DATA - THIS SLIDE CAN BE IGNORED",15,RED,bold=True,first=True,sa=10)
        para(tf,d.get("summary") or "Not assessed. Provide the access on the Basis-of-audit slide to populate it.",13,MGRAY,spc=1.15)
        continue
    chunks=[fs[k:k+4] for k in range(0,len(fs),4)]
    for ci,chunk in enumerate(chunks):
        s=slide(); H(s,"Dimension "+str(i+1)+("  ·  cont." if ci else ""),d.get("name",""), d.get("summary") if ci==0 else None)
        sc=d.get("score","-")
        chip(s,Inches(9.5),Inches(0.62),Inches(1.0),str(sc)+"/5",score_color(sc) if sc!="-" else MGRAY)
        if d.get("confidence"): chip(s,Inches(10.65),Inches(0.62),Inches(2.05),"Confidence: "+str(d["confidence"]),MGRAY)
        y=Inches(2.25) if ci==0 and d.get("summary") else Inches(1.6)
        rows=[[h[0] for h in HDR]]
        for f in chunk:
            rows.append([f,None,None])  # placeholder; fill below
        t=table(s,ML,y,CW,CW3,rows,header_h=0.42,row_h=1.18 if (ci==0 and d.get("summary")) else 1.25)
        for ci2,h in enumerate(HDR): cell(t.cell(0,ci2),[(h[0],10,WHITE,True)],INK)
        for ri,f in enumerate(chunk,start=1):
            fl=WHITE if ri%2 else LGRAY
            cell(t.cell(ri,0),[ (f.get("title",""),11,INK,True), (f.get("current_state",""),9.5,MGRAY,False) ],fl,va=MSO_ANCHOR.TOP)
            cell(t.cell(ri,1),[ (f.get("recommendation",""),10,RGBColor(0x2A,0x2A,0x30),False) ],fl,va=MSO_ANCHOR.TOP)
            cell(t.cell(ri,2),[
                ("Priority:  "+str(f.get("priority","-")),9.5,pcolor(f.get("priority")),True),
                ("Impact:  "+str(f.get("impact","-")),9.5,INK,False),
                ("Effort:  "+str(f.get("effort","-")),9.5,MGRAY,False),
            ],fl,va=MSO_ANCHOR.TOP)

# ============================== 10. NEW RELIC SLOWEST TRANSACTIONS ==============================
tx=data.get("transactions")
if tx:
    s=slide(); H(s,"07  Performance evidence","Page performance (New Relic, production)",
      "Fastly cache serves ~80-85% of catalog in <0.3s. Apdex reflects uncached/first-hit + uncacheable pages (Cart/Checkout). Lower Apdex = more frustrated users.")
    rows=[["PAGE / TRANSACTION","MONTHLY","AVG LOAD","APDEX"]]+[[t.get("name",""),t.get("visits",t.get("count","-")),t.get("avg",t.get("p50","-")),str(t.get("apdex",t.get("p95","-")))] for t in tx]
    tb=table(s,ML,Inches(2.4),CW,[Inches(5.9),Inches(1.9),Inches(2.1),Inches(2.2)],rows,row_h=0.46,header_h=0.44)
    for ri,row in enumerate(rows):
        for ci,val in enumerate(row):
            if ri==0: cell(tb.cell(ri,ci),[(val,10,WHITE,True,PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER)],INK)
            else:
                fl=WHITE if ri%2 else LGRAY
                if ci==0: cell(tb.cell(ri,ci),[(val,10.5,INK,True)],fl)
                elif ci==3: cell(tb.cell(ri,ci),[(val,10.5,apdex_color(val),True,PP_ALIGN.CENTER)],fl)
                else: cell(tb.cell(ri,ci),[(val,10.5,MGRAY,False,PP_ALIGN.CENTER)],fl)

# ============================== 11. ENGINEERING PLAN (perf + code quality) ==============================
ep=data.get("engineering_plan")
if ep:
    s=slide(); H(s,"08  Engineering plan","Performance & code-quality tasks")
    n=len(ep[:3]); cw=Inches((12.1-(n-1)*0.25)/max(1,n))
    for i,sec in enumerate(ep[:3]):
        x=ML+i*(cw+Inches(0.25))
        rect(s,x,Inches(1.85),cw,Inches(4.8),fill=LGRAY); rect(s,x,Inches(1.85),cw,Inches(0.6),fill=INK)
        para(tbox(s,x+Inches(0.22),Inches(1.96),cw-Inches(0.4),Inches(0.4),MSO_ANCHOR.MIDDLE),sec.get("area",""),12.5,WHITE,bold=True,first=True)
        tf=tbox(s,x+Inches(0.28),Inches(2.65),cw-Inches(0.5),Inches(3.8))
        for j,it in enumerate(sec.get("items",[])): para(tf,it,11,INK,first=(j==0),bullet=True,sa=8,spc=1.05)

# ============================== 9b. EXPECTED OUTCOMES ==============================
po=data.get("perf_outcomes")
if po:
    s=slide(); H(s,"09  Expected outcomes","Projected result, the tasks that deliver it, and how to verify",
      "CURRENT measured in New Relic APM (production) + staging Magento profiler / Lighthouse. EXPECTED is a realistic range, achieved by the listed tasks, and confirmed in the named tool.")
    rows=[["METRIC / PAGE","CURRENT","EXPECTED","KEY TASKS (what delivers it)","VERIFY IN"]]+[[r.get("metric",""),r.get("current","-"),r.get("target","-"),r.get("tasks",r.get("note","")),r.get("verify","-")] for r in po[:8]]
    tb=table(s,ML,Inches(2.35),CW,[Inches(2.3),Inches(1.6),Inches(1.6),Inches(4.1),Inches(2.5)],rows,row_h=0.54,header_h=0.46)
    for ri,row in enumerate(rows):
        for ci,val in enumerate(row):
            if ri==0: cell(tb.cell(ri,ci),[(val,9,WHITE,True,PP_ALIGN.LEFT if ci in (0,3,4) else PP_ALIGN.CENTER)],INK)
            else:
                fl=WHITE if ri%2 else LGRAY
                if ci==0: cell(tb.cell(ri,ci),[(val,9.5,INK,True)],fl,va=MSO_ANCHOR.TOP)
                elif ci==1: cell(tb.cell(ri,ci),[(val,9.5,RED,True,PP_ALIGN.CENTER)],fl,va=MSO_ANCHOR.TOP)
                elif ci==2: cell(tb.cell(ri,ci),[(val,9.5,GREEN,True,PP_ALIGN.CENTER)],fl,va=MSO_ANCHOR.TOP)
                elif ci==3: cell(tb.cell(ri,ci),[(val,8.5,INK,False)],fl,va=MSO_ANCHOR.TOP)
                else: cell(tb.cell(ri,ci),[(val,8.5,MGRAY,False)],fl,va=MSO_ANCHOR.TOP)

# ============================== 12. ACTION PLAN ==============================
s=slide(); H(s,"10  Action plan","Quick wins vs roadmap")
qw=data.get("quick_wins",[]); rm=data.get("roadmap",[])
rect(s,ML,Inches(1.7),Inches(5.85),Inches(5.0),fill=LGRAY); rect(s,ML,Inches(1.7),Inches(5.85),Inches(0.55),fill=GREEN)
para(tbox(s,ML+Inches(0.25),Inches(1.8),Inches(5.4),Inches(0.4),MSO_ANCHOR.MIDDLE),"QUICK WINS",13,WHITE,bold=True,first=True)
tf=tbox(s,ML+Inches(0.3),Inches(2.45),Inches(5.3),Inches(4.0))
for j,it in enumerate(qw): para(tf,str(it),12,INK,first=(j==0),bullet=True,sa=10,spc=1.05)
rx=Inches(6.85)
rect(s,rx,Inches(1.7),Inches(5.85),Inches(5.0),fill=LGRAY); rect(s,rx,Inches(1.7),Inches(5.85),Inches(0.55),fill=RED)
para(tbox(s,rx+Inches(0.25),Inches(1.8),Inches(5.4),Inches(0.4),MSO_ANCHOR.MIDDLE),"ROADMAP",13,WHITE,bold=True,first=True)
tf=tbox(s,rx+Inches(0.3),Inches(2.45),Inches(5.3),Inches(4.0))
for j,ph in enumerate(rm):
    if isinstance(ph,dict):
        para(tf,ph.get("phase","")+("   ("+ph.get('effort')+")" if ph.get("effort") else ""),12,INK,bold=True,first=(j==0),sa=2,bullet=True)
        for it in ph.get("items",[]): para(tf,it,10.5,MGRAY,sa=2)
        para(tf," ",6,MGRAY,sa=2)
    else: para(tf,str(ph),12,INK,first=(j==0),bullet=True,sa=8)

# ============================== 13. BENCHMARK ==============================
bm=data.get("benchmark")
if bm and bm.get("rows"):
    s=slide(); H(s,"11  Competitive benchmark",bm.get("category","Category peers"), bm.get("note"))
    peers=bm.get("peers",[]); selfn=CLIENT.split(" - ")[0].split("(")[0].strip()[:14]
    hdr=["FEATURE / EXPERIENCE",selfn]+peers
    pw=Inches(round((12.1-3.8-1.9)/max(1,len(peers)),2)); widths=[Inches(3.8),Inches(1.9)]+[pw]*len(peers)
    rows=[hdr]+[[r.get("feature","")]+[r.get("self","")]+[str(x) for x in r.get("peers",[])] for r in bm["rows"]]
    tb=table(s,ML,Inches(2.5),CW,widths,rows,row_h=0.38,header_h=0.42)
    for ri,row in enumerate(rows):
        for ci,val in enumerate(row):
            if ri==0: cell(tb.cell(ri,ci),[(val,9.5,WHITE,True,PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER)],INK)
            else:
                fl=WHITE if ri%2 else LGRAY
                if ci==0: cell(tb.cell(ri,ci),[(val,9.5,INK,True)],fl)
                elif ci==1:
                    sv=str(val).strip().lower(); col=GREEN if sv.startswith(("yes","strong","best","good")) else (RED if sv.startswith(("no","broken","weak","absent","missing","none","off")) else INK)
                    cell(tb.cell(ri,ci),[(val,9.5,col,True,PP_ALIGN.CENTER)],fl)
                else: cell(tb.cell(ri,ci),[(val,9,MGRAY,False,PP_ALIGN.CENTER)],fl)

# ============================== 12. MATRIX (moved to bottom) ==============================
s=slide(); H(s,"12  Prioritisation","Impact vs Effort")
m=data.get("matrix")
if not m:
    m={"do_first":[],"plan_next":[],"quick_wins":[],"deprioritize":[]}
    for d in dims:
        for f in d.get("findings",[]):
            imp=str(f.get("impact","")).lower(); eff=str(f.get("effort","")).lower(); lbl=f.get("id") or f.get("title","")
            if imp=="high" and eff=="low": m["do_first"].append(lbl)
            elif imp=="high": m["plan_next"].append(lbl)
            elif eff=="low": m["quick_wins"].append(lbl)
            else: m["deprioritize"].append(lbl)
quad=[("DO FIRST","high impact / low effort","do_first",GREEN,ML,Inches(1.85)),
      ("PLAN NEXT","high impact / high effort","plan_next",AMBER,Inches(6.78),Inches(1.85)),
      ("QUICK WINS","low impact / low effort","quick_wins",INK,ML,Inches(4.5)),
      ("DEPRIORITISE","low impact / high effort","deprioritize",MGRAY,Inches(6.78),Inches(4.5))]
for lab,subl,key,col,x,y in quad:
    rect(s,x,y,Inches(5.95),Inches(2.5),fill=LGRAY); rect(s,x,y,Inches(5.95),Inches(0.5),fill=col)
    para(tbox(s,x+Inches(0.22),y+Inches(0.08),Inches(5.5),Inches(0.34),MSO_ANCHOR.MIDDLE),lab+"   -   "+subl,11.5,WHITE,bold=True,first=True)
    tf=tbox(s,x+Inches(0.3),y+Inches(0.62),Inches(5.4),Inches(1.8)); its=m.get(key,[])
    if not its: para(tf,"-",11,MGRAY,first=True)
    for j,it in enumerate(its[:7]): para(tf,str(it),11,INK,first=(j==0),bullet=True,sa=3)

# ============================== 14. CLOSE ==============================
s=slide(); rect(s,0,0,SW,SH,fill=INK); rect(s,ML,Inches(2.85),Inches(0.85),Inches(0.09),fill=RED)
para(tbox(s,ML,Inches(1.9),CW,Inches(0.9)),"Expected impact & next steps",36,WHITE,bold=True,first=True)
tf=tbox(s,ML,Inches(3.2),CW,Inches(3.2))
ns=data.get("next_steps")
if isinstance(ns,list):
    for j,it in enumerate(ns): para(tf,str(it),15,RGBColor(0xD0,0xD0,0xD6),first=(j==0),bullet=True,sa=12,spc=1.1)
elif ns: para(tf,str(ns),16,RGBColor(0xD0,0xD0,0xD6),first=True)

prs.save(out)
print("Saved:",out,"| slides:",len(prs.slides._sldIdLst))
